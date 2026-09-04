from __future__ import annotations

import re
from collections.abc import Callable
from pathlib import Path

from rag_docs.models import DocumentCandidate, ExtractedUnit

# Bump when the extraction logic changes what text a document yields for the
# same bytes: it is one of the components IndexFingerprint covers (RULE-004).
EXTRACTOR_VERSION = "extractors-v1"


class ExtractionError(RuntimeError):
    pass


class UnsupportedDocumentError(ExtractionError):
    pass


def _clean(text: str) -> str:
    return re.sub(r"[ \t]+", " ", re.sub(r"\r\n?", "\n", text)).strip()


def _extract_pdf(path: Path) -> list[ExtractedUnit]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return [
        ExtractedUnit(text=text, locator={"page": index})
        for index, page in enumerate(reader.pages, start=1)
        if (text := _clean(page.extract_text() or ""))
    ]


def _extract_docx(path: Path) -> list[ExtractedUnit]:
    from docx import Document

    document = Document(str(path))
    units: list[ExtractedUnit] = []
    section: str | None = None
    buffer: list[str] = []

    def flush() -> None:
        text = _clean("\n".join(buffer))
        if text:
            units.append(ExtractedUnit(text=text, section=section))
        buffer.clear()

    for paragraph in document.paragraphs:
        text = _clean(paragraph.text)
        if not text:
            continue
        if paragraph.style and paragraph.style.name.lower().startswith("heading"):
            flush()
            section = text
        else:
            buffer.append(text)
    flush()

    for number, table in enumerate(document.tables, start=1):
        rows = ["\t".join(_clean(cell.text) for cell in row.cells) for row in table.rows]
        text = _clean("\n".join(rows))
        if text:
            units.append(
                ExtractedUnit(text=text, locator={"table": number}, section=section)
            )
    return units


def _extract_pptx(path: Path) -> list[ExtractedUnit]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    units: list[ExtractedUnit] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        title: str | None = None
        if slide.shapes.title:
            title = _clean(slide.shapes.title.text)
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _clean(shape.text)
                if text and text != title:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                rows = [
                    "\t".join(_clean(cell.text) for cell in row.cells)
                    for row in shape.table.rows
                ]
                parts.extend(row for row in rows if row.strip())
        text = _clean("\n".join(parts))
        if title:
            text = _clean(f"{title}\n{text}")
        if text:
            units.append(
                ExtractedUnit(
                    text=text, locator={"slide": slide_number}, section=title or None
                )
            )
    return units


def _extract_xlsx(path: Path) -> list[ExtractedUnit]:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    workbook = load_workbook(path, read_only=True, data_only=True)
    units: list[ExtractedUnit] = []
    try:
        for sheet in workbook.worksheets:
            block: list[str] = []
            start_row: int | None = None
            last_row = 0
            max_column = 1

            for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                values = ["" if value is None else str(value) for value in row]
                while values and not values[-1]:
                    values.pop()
                if not any(values):
                    if block and start_row is not None:
                        units.append(
                            ExtractedUnit(
                                text="\n".join(block),
                                locator={
                                    "sheet": sheet.title,
                                    "cell_range": (
                                        f"A{start_row}:"
                                        f"{get_column_letter(max_column)}{last_row}"
                                    ),
                                },
                                section=sheet.title,
                            )
                        )
                    block = []
                    start_row = None
                    continue
                if start_row is None:
                    start_row = row_number
                last_row = row_number
                max_column = max(max_column, len(values))
                block.append("\t".join(values))
                if len(block) >= 50:
                    units.append(
                        ExtractedUnit(
                            text="\n".join(block),
                            locator={
                                "sheet": sheet.title,
                                "cell_range": (
                                    f"A{start_row}:"
                                    f"{get_column_letter(max_column)}{last_row}"
                                ),
                            },
                            section=sheet.title,
                        )
                    )
                    block = []
                    start_row = None
            if block and start_row is not None:
                units.append(
                    ExtractedUnit(
                        text="\n".join(block),
                        locator={
                            "sheet": sheet.title,
                            "cell_range": (
                                f"A{start_row}:"
                                f"{get_column_letter(max_column)}{last_row}"
                            ),
                        },
                        section=sheet.title,
                    )
                )
    finally:
        workbook.close()
    return units


def _extract_text(path: Path) -> list[ExtractedUnit]:
    text = _clean(path.read_text(encoding="utf-8-sig", errors="replace"))
    return [ExtractedUnit(text=text)] if text else []


def _extract_markdown(path: Path) -> list[ExtractedUnit]:
    content = path.read_text(encoding="utf-8-sig", errors="replace")
    units: list[ExtractedUnit] = []
    section: str | None = None
    buffer: list[str] = []
    for line in content.splitlines():
        heading = re.match(r"^#{1,6}\s+(.+?)\s*$", line)
        if heading:
            text = _clean("\n".join(buffer))
            if text:
                units.append(ExtractedUnit(text=text, section=section))
            section = heading.group(1).strip()
            buffer = []
        else:
            buffer.append(line)
    text = _clean("\n".join(buffer))
    if text:
        units.append(ExtractedUnit(text=text, section=section))
    return units


EXTRACTORS: dict[str, Callable[[Path], list[ExtractedUnit]]] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".txt": _extract_text,
    ".md": _extract_markdown,
}


def extract_document(candidate: DocumentCandidate) -> list[ExtractedUnit]:
    extractor = EXTRACTORS.get(candidate.path.suffix.casefold())
    if extractor is None:
        raise UnsupportedDocumentError(f"Formato no soportado: {candidate.path.suffix}")
    try:
        return extractor(candidate.path)
    except UnsupportedDocumentError:
        raise
    except Exception as exc:
        raise ExtractionError(f"No se pudo extraer {candidate.relative_path}: {exc}") from exc
