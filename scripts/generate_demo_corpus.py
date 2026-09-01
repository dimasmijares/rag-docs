"""Genera un corpus didáctico multiformato reproducible y su manifiesto."""

from __future__ import annotations

import argparse
import hashlib
import io
import re
import tempfile
import zipfile
from datetime import UTC, datetime
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_ROOT = PROJECT_ROOT / "examples" / "corpus" / "demo"
CORPUS_VERSION = "0.2.0"
SCHEMA_VERSION = "1.0"
FIXED_DATETIME = datetime(2026, 1, 1, tzinfo=UTC)
FIXED_ZIP_TIME = (1980, 1, 1, 0, 0, 0)
MANIFEST_NAME = "manifest.sha256"
FIXTURE_PATHS = (
    "architecture/atlas-routing.pptx",
    "architecture/flujo-datos.pptx",
    "components/orquestador.docx",
    "continuity/atlas-guide.md",
    "continuity/atlas-procedure.docx",
    "etl/catalogo-etl.md",
    "inventory/atlas-services.xlsx",
    "inventory/inventario-procesos.xlsx",
    "operations/atlas-oncall.txt",
    "operations/incidencias.txt",
    "runbooks/atlas-rollback.pdf",
    "runbooks/recuperacion-clientes.pdf",
)


def _write_text(root: Path, relative_path: str, content: str) -> None:
    path = root / relative_path
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content.rstrip() + "\n", encoding="utf-8", newline="\n")


def _canonicalize_ooxml(buffer: io.BytesIO) -> bytes:
    """Ordena y normaliza los miembros ZIP para obtener OOXML byte a byte estable."""
    source = zipfile.ZipFile(io.BytesIO(buffer.getvalue()), "r")
    output = io.BytesIO()
    with source, zipfile.ZipFile(
        output, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
    ) as target:
        for name in sorted(source.namelist()):
            data = source.read(name)
            if name == "docProps/core.xml":
                for field in (b"created", b"modified"):
                    data = re.sub(
                        rb"(<dcterms:" + field + rb"[^>]*>).*?(</dcterms:" + field + rb">)",
                        rb"\g<1>2026-01-01T00:00:00Z\g<2>",
                        data,
                    )
            info = zipfile.ZipInfo(name, FIXED_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.create_system = 0
            info.external_attr = 0
            target.writestr(
                info,
                data,
                compress_type=zipfile.ZIP_DEFLATED,
                compresslevel=9,
            )
    return output.getvalue()


def _save_ooxml(root: Path, relative_path: str, artifact: object) -> None:
    path = root / relative_path
    path.parent.mkdir(parents=True, exist_ok=True)
    buffer = io.BytesIO()
    artifact.save(buffer)  # type: ignore[attr-defined]
    path.write_bytes(_canonicalize_ooxml(buffer))


def create_markdown(root: Path) -> None:
    _write_text(
        root,
        "etl/catalogo-etl.md",
        """# Catálogo de ETL

## ETL_CLIENTES_DIARIA

La ETL `ETL_CLIENTES_DIARIA` carga `CLIENTE_MAESTRO` desde `CRM_ORACLE` cada día a las
02:00. Su definición técnica está en `procesos/clientes/etl_clientes_diaria`.

## ETL_CONTRATOS_SEMANAL

La ETL `ETL_CONTRATOS_SEMANAL` consolida contratos activos cada domingo a las 06:00.

## Gobierno del orquestador

El equipo Plataforma de Datos mantiene `ORQ_DATAOPS` y aprueba sus cambios de calendario.
""",
    )


def create_text(root: Path) -> None:
    _write_text(
        root,
        "operations/incidencias.txt",
        """RUNBOOK DE INCIDENCIAS

Si ETL_CLIENTES_DIARIA falla, revisar primero la conexión CRM_ORACLE y después el fichero
de rechazo en /operaciones/rechazos/clientes.
El equipo DataOps atiende la primera intervención; Plataforma de Datos recibe la escalada.
El código de diagnóstico para un timeout de origen es SRC_TIMEOUT_17.
""",
    )


def create_atlas_markdown(root: Path) -> None:
    _write_text(
        root,
        "continuity/atlas-guide.md",
        """# Continuidad de Atlas

## Restauración Atlas

La restauración de `ATLAS_LEDGER` usa `SNAPSHOT_BLUE` y requiere aprobación de `RISK_OPS`.
""",
    )


def create_atlas_text(root: Path) -> None:
    _write_text(
        root,
        "operations/atlas-oncall.txt",
        """GUARDIA DE ATLAS

La guardia revisa primero METRIC_LAG_42. Tras 15 minutos escala a DATA_RELIABILITY.
""",
    )


def create_docx(root: Path) -> None:
    document = Document()
    document.core_properties.title = "Orquestador de cargas sintético"
    document.core_properties.author = "rag-docs-team"
    document.core_properties.last_modified_by = "rag-docs-team"
    document.core_properties.created = FIXED_DATETIME
    document.core_properties.modified = FIXED_DATETIME
    document.add_heading("Orquestador de cargas", level=1)
    document.add_paragraph(
        "ORQ_DATAOPS controla las dependencias de las ETL nocturnas y pertenece a Plataforma "
        "de Datos."
    )
    document.add_heading("Reintentos", level=2)
    document.add_paragraph(
        "Cada ejecución admite dos reintentos automáticos separados por cinco minutos."
    )
    document.add_heading("Monitorización", level=2)
    document.add_paragraph(
        "Las ejecuciones se consultan en CONTROL_M y las alertas críticas usan el código "
        "ORQ_CRIT_9."
    )
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Severidad"
    table.rows[0].cells[1].text = "Canal"
    row = table.add_row().cells
    row[0].text = "Crítica"
    row[1].text = "DATAOPS_GUARDIA"
    _save_ooxml(root, "components/orquestador.docx", document)


def create_atlas_docx(root: Path) -> None:
    document = Document()
    document.core_properties.title = "Procedimiento sintético de Atlas"
    document.core_properties.author = "rag-docs-team"
    document.core_properties.last_modified_by = "rag-docs-team"
    document.core_properties.created = FIXED_DATETIME
    document.core_properties.modified = FIXED_DATETIME
    document.add_heading("Procedimiento de continuidad", level=1)
    document.add_heading("Restauración Atlas", level=2)
    document.add_paragraph(
        "La restauración de ATLAS_LEDGER usa SNAPSHOT_BLUE y requiere aprobación de RISK_OPS."
    )
    _save_ooxml(root, "continuity/atlas-procedure.docx", document)


def create_pptx(root: Path) -> None:
    presentation = Presentation()
    presentation.core_properties.title = "Flujo de datos sintético"
    presentation.core_properties.author = "rag-docs-team"
    presentation.core_properties.last_modified_by = "rag-docs-team"
    presentation.core_properties.created = FIXED_DATETIME
    presentation.core_properties.modified = FIXED_DATETIME
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Flujo de datos de clientes"
    slide.placeholders[1].text = (
        "CRM_ORACLE -> ETL_CLIENTES_DIARIA -> CLIENTE_MAESTRO -> reporting\n"
        "La validación de disponibilidad se ejecuta a las 04:30."
    )
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Controles de publicación"
    slide.placeholders[1].text = (
        "El control PUB_CLIENTES_21 compara el recuento de origen y destino antes de publicar."
    )
    _save_ooxml(root, "architecture/flujo-datos.pptx", presentation)


def create_atlas_pptx(root: Path) -> None:
    presentation = Presentation()
    presentation.core_properties.title = "Ruta Atlas sintética"
    presentation.core_properties.author = "rag-docs-team"
    presentation.core_properties.last_modified_by = "rag-docs-team"
    presentation.core_properties.created = FIXED_DATETIME
    presentation.core_properties.modified = FIXED_DATETIME
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Ruta Atlas"
    slide.placeholders[1].text = "KAFKA_ATLAS -> NORMALIZE_V2 -> LEDGER_CURATED"
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Servicio Atlas"
    slide.placeholders[1].text = "PlatformOps opera la ruta con un SLA de 12 minutos."
    _save_ooxml(root, "architecture/atlas-routing.pptx", presentation)


def create_xlsx(root: Path) -> None:
    workbook = Workbook()
    workbook.properties.title = "Inventario de procesos sintético"
    workbook.properties.creator = "rag-docs-team"
    workbook.properties.lastModifiedBy = "rag-docs-team"
    workbook.properties.created = FIXED_DATETIME.replace(tzinfo=None)
    workbook.properties.modified = FIXED_DATETIME.replace(tzinfo=None)
    sheet = workbook.active
    sheet.title = "Procesos"
    sheet.append(["Proceso", "Responsable", "SLA"])
    sheet.append(["ETL_CLIENTES_DIARIA", "DataOps", "03:00"])
    sheet.append(["ETL_CONTRATOS_SEMANAL", "Backoffice", "Lunes 06:00"])
    sheet.append(["PUB_CLIENTES_21", "Calidad de Datos", "04:45"])
    sheet.freeze_panes = "A2"
    sheet.column_dimensions["A"].width = 28
    sheet.column_dimensions["B"].width = 20
    sheet.column_dimensions["C"].width = 18
    for cell in sheet[1]:
        cell.font = Font(bold=True)
    contacts = workbook.create_sheet("Contactos")
    contacts.append(["Equipo", "Buzón sintético"])
    contacts.append(["DataOps", "dataops@example.invalid"])
    contacts.append(["Backoffice", "backoffice@example.invalid"])
    contacts.freeze_panes = "A2"
    contacts.column_dimensions["A"].width = 18
    contacts.column_dimensions["B"].width = 34
    for cell in contacts[1]:
        cell.font = Font(bold=True)
    _save_ooxml(root, "inventory/inventario-procesos.xlsx", workbook)


def create_atlas_xlsx(root: Path) -> None:
    workbook = Workbook()
    workbook.properties.title = "Servicios Atlas sintéticos"
    workbook.properties.creator = "rag-docs-team"
    workbook.properties.lastModifiedBy = "rag-docs-team"
    workbook.properties.created = FIXED_DATETIME.replace(tzinfo=None)
    workbook.properties.modified = FIXED_DATETIME.replace(tzinfo=None)
    sheet = workbook.active
    sheet.title = "Servicios"
    sheet.append(["Servicio", "Responsable", "SLO"])
    sheet.append(["ATLAS_GATEWAY", "IntegrationOps", "99.95%"])
    sheet.append(["ATLAS_ARCHIVE", "RecordsOps", "99.90%"])
    sheet.freeze_panes = "A2"
    sheet.column_dimensions["A"].width = 22
    sheet.column_dimensions["B"].width = 20
    sheet.column_dimensions["C"].width = 14
    for cell in sheet[1]:
        cell.font = Font(bold=True)
    _save_ooxml(root, "inventory/atlas-services.xlsx", workbook)


def _pdf_page(writer: PdfWriter, text: str) -> None:
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    page[NameObject("/Resources")] = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})}
    )
    stream = DecodedStreamObject()
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream.set_data(f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode("ascii"))
    page[NameObject("/Contents")] = writer._add_object(stream)


def create_pdf(root: Path) -> None:
    writer = PdfWriter()
    writer.add_metadata(
        {
            "/Title": "Runbook sintético de recuperación",
            "/Author": "rag-docs-team",
            "/CreationDate": "D:20260101000000Z",
            "/ModDate": "D:20260101000000Z",
        }
    )
    _pdf_page(
        writer,
        "Reintentar ETL_CLIENTES_DIARIA desde CONTROL_M tras validar CRM_ORACLE.",
    )
    _pdf_page(
        writer,
        "La comprobacion posterior se ejecuta a las 04:30 y registra RECOVERY_OK_42.",
    )
    path = root / "runbooks" / "recuperacion-clientes.pdf"
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("wb") as output:
        writer.write(output)


def create_atlas_pdf(root: Path) -> None:
    writer = PdfWriter()
    writer.add_metadata(
        {
            "/Title": "Rollback sintético de Atlas",
            "/Author": "rag-docs-team",
            "/CreationDate": "D:20260101000000Z",
            "/ModDate": "D:20260101000000Z",
        }
    )
    _pdf_page(writer, "Pausar ATLAS_WRITER y restaurar SNAPSHOT_BLUE.")
    _pdf_page(writer, "Reanudar solo despues de CHECKSUM_GREEN.")
    path = root / "runbooks" / "atlas-rollback.pdf"
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("wb") as output:
        writer.write(output)


def _fixture_paths(root: Path) -> list[Path]:
    return [root / relative_path for relative_path in FIXTURE_PATHS]


def create_manifest(root: Path) -> None:
    lines = [
        f"# corpus_version: {CORPUS_VERSION}",
        f"# schema_version: {SCHEMA_VERSION}",
    ]
    for path in _fixture_paths(root):
        digest = hashlib.sha256(path.read_bytes()).hexdigest()
        lines.append(f"{digest}  {path.relative_to(root).as_posix()}")
    _write_text(root, MANIFEST_NAME, "\n".join(lines))


def generate_corpus(root: Path) -> None:
    root.mkdir(parents=True, exist_ok=True)
    create_markdown(root)
    create_text(root)
    create_atlas_markdown(root)
    create_atlas_text(root)
    create_docx(root)
    create_atlas_docx(root)
    create_pptx(root)
    create_atlas_pptx(root)
    create_xlsx(root)
    create_atlas_xlsx(root)
    create_pdf(root)
    create_atlas_pdf(root)
    create_manifest(root)


def _compare_trees(expected: Path, actual: Path) -> list[str]:
    expected_files = {
        path.relative_to(expected).as_posix(): path.read_bytes()
        for path in expected.rglob("*")
        if path.is_file()
    }
    actual_files = {
        path.relative_to(actual).as_posix(): path.read_bytes()
        for path in actual.rglob("*")
        if path.is_file()
    }
    differences = sorted(set(expected_files) ^ set(actual_files))
    differences.extend(
        path
        for path in sorted(set(expected_files) & set(actual_files))
        if expected_files[path] != actual_files[path]
    )
    return differences


def check_corpus(root: Path) -> None:
    with tempfile.TemporaryDirectory(prefix="rag-docs-corpus-") as directory:
        generated = Path(directory)
        generate_corpus(generated)
        differences = _compare_trees(root, generated)
    if differences:
        raise SystemExit("Corpus no canónico: " + ", ".join(differences))
    print(f"Corpus canónico verificado en {root}")


def run() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_ROOT)
    parser.add_argument("--check", action="store_true")
    args = parser.parse_args()
    root = args.output_dir.resolve()
    if args.check:
        check_corpus(root)
        return
    generate_corpus(root)
    print(f"Corpus generado en {root}")


if __name__ == "__main__":
    run()
